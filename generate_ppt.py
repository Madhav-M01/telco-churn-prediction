"""Generate Telco Churn Prediction PPT using the project template."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy

TEMPLATE = 'template (1).pptx'
OUTPUT   = 'Churn_Prediction_Presentation_v2.pptx'

# ── Colour palette ────────────────────────────────────────────────────────────
C_DARK_BLUE  = RGBColor(0x1F, 0x49, 0x7D)
C_MID_BLUE   = RGBColor(0x2E, 0x74, 0xB5)
C_LIGHT_BLUE = RGBColor(0xBD, 0xD7, 0xEE)
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK      = RGBColor(0x00, 0x00, 0x00)
C_GRAY       = RGBColor(0x40, 0x40, 0x40)
C_DARK_GRAY  = RGBColor(0x26, 0x26, 0x26)
C_GREEN      = RGBColor(0x37, 0x86, 0x35)
C_ORANGE     = RGBColor(0xED, 0x7D, 0x31)
C_RED        = RGBColor(0xC0, 0x00, 0x00)
C_PURPLE     = RGBColor(0x5C, 0x27, 0x8A)

# ── Load template ─────────────────────────────────────────────────────────────
prs          = Presentation(TEMPLATE)
BLANK_LAYOUT = prs.slide_layouts[6]
header_src   = prs.slides[1]   # slide 2 has the logo header picture


# ── Helper: copy header image to any new slide ────────────────────────────────
def copy_header(dest_slide):
    NS_A   = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    NS_R   = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    NS_IMG = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/image'
    for shape in header_src.shapes:
        if shape.shape_type == 13:                      # Picture
            sp_copy = copy.deepcopy(shape._element)
            blip = sp_copy.find(f'.//{{{NS_A}}}blip')
            if blip is not None:
                old_rId = blip.get(f'{{{NS_R}}}embed')
                if old_rId:
                    img_part = header_src.part.related_part(old_rId)
                    new_rId  = dest_slide.part.relate_to(img_part, NS_IMG)
                    blip.set(f'{{{NS_R}}}embed', new_rId)
            dest_slide.shapes._spTree.append(sp_copy)
            break


# ── Helper: title bar + divider line ─────────────────────────────────────────
def add_title_bar(slide, text):
    box = slide.shapes.add_textbox(Inches(0.4), Inches(1.25), Inches(12.53), Inches(0.48))
    tf  = box.text_frame
    p   = tf.paragraphs[0]
    r   = p.add_run()
    r.text           = text
    r.font.size      = Pt(24)
    r.font.bold      = True
    r.font.color.rgb = C_DARK_BLUE
    div = slide.shapes.add_shape(1, Inches(0.4), Inches(1.76), Inches(12.53), Inches(0.045))
    div.fill.solid()
    div.fill.fore_color.rgb = C_MID_BLUE
    div.line.fill.background()


# ── Helper: simple text box ───────────────────────────────────────────────────
def txb(slide, left, top, width, height, text, size=14, bold=False,
        color=C_GRAY, align=PP_ALIGN.LEFT, italic=False, wrap=True):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    r   = p.add_run()
    r.text           = text
    r.font.size      = Pt(size)
    r.font.bold      = bold
    r.font.italic    = italic
    r.font.color.rgb = color
    return box


# ── Helper: multi-line textbox with bullets ───────────────────────────────────
def bullets(slide, left, top, width, height, lines):
    """lines = list of (text, size, bold, color, use_bullet)"""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf  = box.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            item = (item, 14, False, C_GRAY, True)
        text, size, bold, color, use_bullet = item
        p     = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r     = p.add_run()
        r.text           = ('  >>  ' if use_bullet else '') + text
        r.font.size      = Pt(size)
        r.font.bold      = bold
        r.font.color.rgb = color
    return box


# ── Helper: styled table ──────────────────────────────────────────────────────
def add_table(slide, left, top, width, rows_data, col_widths,
              hdr_color=C_DARK_BLUE, alt_color=RGBColor(0xF2, 0xF2, 0xF2),
              font_size=12):
    rows  = len(rows_data)
    cols  = len(col_widths)
    row_h = 0.37
    tbl   = slide.shapes.add_table(
        rows, cols,
        Inches(left), Inches(top),
        Inches(width), Inches(rows * row_h)
    ).table
    total_w = sum(col_widths)
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = Inches(width * cw / total_w)
    for ri, row in enumerate(rows_data):
        for ci, cell_text in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text_frame.word_wrap = True
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
            # Get or create run
            r = p.runs[0] if p.runs else p.add_run()
            r.text           = str(cell_text)
            r.font.size      = Pt(font_size)
            r.font.bold      = (ri == 0)
            if ri == 0:
                r.font.color.rgb = C_WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = hdr_color
            elif ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = alt_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_WHITE


# ── Helper: colored metric / info box ────────────────────────────────────────
def info_box(slide, left, top, width, height, label, value,
             bg=C_LIGHT_BLUE, label_color=C_DARK_BLUE, val_color=C_DARK_BLUE):
    rect = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    rect.fill.solid()
    rect.fill.fore_color.rgb = bg
    rect.line.color.rgb = C_MID_BLUE
    tb1  = slide.shapes.add_textbox(
        Inches(left + 0.05), Inches(top + 0.05),
        Inches(width - 0.1),  Inches(height * 0.42))
    p    = tb1.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r    = p.add_run()
    r.text           = label
    r.font.size      = Pt(10)
    r.font.bold      = True
    r.font.color.rgb = label_color
    tb2  = slide.shapes.add_textbox(
        Inches(left + 0.05), Inches(top + height * 0.42),
        Inches(width - 0.1),  Inches(height * 0.58))
    p2   = tb2.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2   = p2.add_run()
    r2.text           = value
    r2.font.size      = Pt(18)
    r2.font.bold      = True
    r2.font.color.rgb = val_color


# ── Helper: new blank content slide ──────────────────────────────────────────
def new_slide(title_text):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    copy_header(slide)
    add_title_bar(slide, title_text)
    return slide


# ── Helper: move slide at old_idx to new_idx (fixes ordering) ────────────────
def move_slide(prs, old_idx, new_idx):
    sldIdLst = prs.slides._sldIdLst
    elem = sldIdLst[old_idx]
    sldIdLst.remove(elem)
    sldIdLst.insert(new_idx, elem)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — TITLE SLIDE  (modify existing slide 0)
# ══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides[0]

for shape in s1.shapes:
    if shape.name == 'Title 1':
        tf = shape.text_frame
        tf.clear()
        p  = tf.add_paragraph()
        r  = p.add_run()
        r.text           = 'Telco Customer Churn Prediction'
        r.font.size      = Pt(34)
        r.font.bold      = True
        r.font.color.rgb = C_DARK_BLUE
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text           = 'End-to-End Production-Grade ML Pipeline\nwith AI-Powered Explainability'
        r2.font.size      = Pt(20)
        r2.font.color.rgb = C_MID_BLUE

    elif shape.name == 'Subtitle 2':
        tf = shape.text_frame
        tf.clear()
        p  = tf.add_paragraph()
        r  = p.add_run()
        r.text           = 'Presented by:'
        r.font.size      = Pt(14)
        r.font.bold      = True
        r.font.color.rgb = C_DARK_GRAY
        for name in ['CHIRAG KUMAR  (UI23EC14)',
                     'MADHAV BAGRI  (UI23EC32)',
                     'PRATEEK SHINGH  (UI23EC44)']:
            px = tf.add_paragraph()
            rx = px.add_run()
            rx.text           = name
            rx.font.size      = Pt(13)
            rx.font.bold      = False
            rx.font.color.rgb = C_GRAY

    elif shape.name == 'TextBox 4':
        tf = shape.text_frame
        tf.clear()
        p  = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r  = p.add_run()
        r.text           = 'Department of Electronics and Communication Engineering  |  IIIT Surat'
        r.font.size      = Pt(12)
        r.font.bold      = True
        r.font.color.rgb = C_WHITE

    elif shape.name == 'TextBox 5':
        tf = shape.text_frame
        tf.clear()
        p  = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r  = p.add_run()
        r.text           = 'Under the Guidance of'
        r.font.size      = Pt(13)
        r.font.bold      = True
        r.font.color.rgb = C_DARK_GRAY
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text           = 'Dr. Sudeep Sharma'
        r2.font.size      = Pt(13)
        r2.font.bold      = False
        r2.font.color.rgb = C_DARK_GRAY
        p3 = tf.add_paragraph()
        p3.alignment = PP_ALIGN.CENTER
        r3 = p3.add_run()
        r3.text           = 'Assistant Professor'
        r3.font.size      = Pt(11)
        r3.font.bold      = False
        r3.font.color.rgb = C_GRAY


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — OUTLINE  (modify existing slide 1)
# ══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides[1]
for shape in s2.shapes:
    if shape.name == 'TextBox 6':
        tf = shape.text_frame
        tf.clear()
        heading = tf.add_paragraph()
        r = heading.add_run()
        r.text           = 'Outline of Presentation'
        r.font.size      = Pt(24)
        r.font.bold      = True
        r.font.color.rgb = C_DARK_BLUE
        outline_items = [
            '1.   Introduction & Motivation',
            '2.   Business Problem Formulation',
            '3.   Dataset Overview',
            '4.   ML Pipeline Architecture',
            '5.   Feature Engineering',
            '6.   Handling Class Imbalance',
            '7.   Model Development (Optuna Tuning)',
            '8.   Ensemble & Probability Calibration',
            '9.   Results & Model Comparison',
            '10.  SHAP Explainability',
            '11.  Deployment - Streamlit Chatbot',
            '12.  Conclusion & Future Work',
        ]
        for item in outline_items:
            p = tf.add_paragraph()
            r = p.add_run()
            r.text           = item
            r.font.size      = Pt(15)
            r.font.color.rgb = C_GRAY


# ══════════════════════════════════════════════════════════════════════════════
#  NOTE: Existing slide 2 (index 2) is a blank slide from the template.
#  It will be moved to the end and used as the Thank You slide.
#  All content slides (s3-s14) are added AFTER the existing 3 slides, so they
#  land at indices 3-14. At the end, we move index 2 to the last position.
# ══════════════════════════════════════════════════════════════════════════════

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — INTRODUCTION
# ══════════════════════════════════════════════════════════════════════════════
s3 = new_slide('Introduction & Motivation')

txb(s3, 0.4, 1.9, 12.5, 0.4, 'What is Customer Churn?',
    size=17, bold=True, color=C_DARK_BLUE)

bullets(s3, 0.4, 2.35, 12.2, 1.4, [
    ('Customer churn = customers discontinuing telecom service', 15, False, C_GRAY, True),
    ('Replacing a lost customer costs 5-25x more than retaining one  (Harvard Business Review)', 14, False, C_GRAY, True),
    ('Early churn prediction enables targeted retention and reduced revenue leakage', 14, False, C_GRAY, True),
])

txb(s3, 0.4, 3.75, 12.5, 0.4, 'Research Objectives',
    size=17, bold=True, color=C_DARK_BLUE)

bullets(s3, 0.4, 4.2, 12.2, 2.0, [
    ('Build a production-ready pipeline with ZERO data leakage', 14, False, C_GRAY, True),
    ('Optimize for business impact - not just accuracy', 14, False, C_GRAY, True),
    ('Reflect asymmetric cost structure: FN costs 10x more than FP', 14, False, C_GRAY, True),
    ('Deliver interpretable, actionable predictions via SHAP + AI chatbot', 14, False, C_GRAY, True),
])

txb(s3, 0.4, 6.3, 12.2, 0.4,
    'Dataset: IBM Telco Customer Churn  |  7,043 customers  |  21 raw features  |  26.5% churn rate',
    size=13, bold=False, color=C_MID_BLUE, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 4 — BUSINESS PROBLEM
# ══════════════════════════════════════════════════════════════════════════════
s4 = new_slide('Business Problem Formulation')

txb(s4, 0.4, 1.9, 6.0, 0.35, 'Cost Parameters',
    size=15, bold=True, color=C_DARK_BLUE)
add_table(s4, 0.4, 2.3, 6.2,
    [['Parameter',                          'Value'],
     ['Avg. Monthly Revenue / Customer',    '$64.76'],
     ['Cost of Missing a Churner (FN)',     '$500'],
     ['Cost of False Alarm (FP)',           '$50'],
     ['FN : FP Cost Ratio',                '10 : 1'],
     ['Retention Success Rate',             '30%'],
     ['Dataset Size',                       '7,043'],
     ['Churn Rate',                         '~26.5%']],
    [3.2, 2.0], font_size=12)

txb(s4, 6.8, 1.9, 5.9, 0.35, 'Optimization Objective',
    size=15, bold=True, color=C_DARK_BLUE)
txb(s4, 6.8, 2.3, 5.9, 0.45,
    'Expected Value (EV) per user:',
    size=13, bold=True, color=C_GRAY)
txb(s4, 6.8, 2.75, 5.9, 0.55,
    '  EV = TP x (0.30 x $500 - $50) - FN x $500',
    size=13, bold=True, color=C_DARK_BLUE)

info_box(s4, 6.8, 3.5,  2.7, 1.1, 'FALSE NEGATIVE', '$500 LOSS',
         bg=RGBColor(0xFF, 0xD7, 0xD7), label_color=C_RED, val_color=C_RED)
info_box(s4, 9.6, 3.5,  2.7, 1.1, 'FALSE POSITIVE', '$50 LOSS',
         bg=RGBColor(0xFF, 0xF3, 0xCD), label_color=C_ORANGE, val_color=C_ORANGE)
info_box(s4, 6.8, 4.75, 5.5, 1.0, 'COST RATIO', '10 : 1  =>  Maximize RECALL',
         bg=C_LIGHT_BLUE, label_color=C_DARK_BLUE, val_color=C_DARK_BLUE)

txb(s4, 0.4, 6.35, 12.5, 0.5,
    'Key Insight: Missing a real churner is 10x more costly than a false alarm.'
    ' Model heavily optimizes Recall over Precision.',
    size=12, italic=True, color=C_MID_BLUE)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 5 — DATASET OVERVIEW
# ══════════════════════════════════════════════════════════════════════════════
s5 = new_slide('Dataset Overview')

stats = [('TOTAL RECORDS', '7,043'), ('RAW FEATURES', '21'),
         ('CHURN RATE', '26.5%'),    ('CLASS RATIO', '73 : 27')]
for i, (lbl, val) in enumerate(stats):
    info_box(s5, 0.4 + i * 3.0, 1.9, 2.8, 1.1, lbl, val, bg=C_LIGHT_BLUE)

txb(s5, 0.4, 3.15, 12.5, 0.38, 'Feature Categories',
    size=15, bold=True, color=C_DARK_BLUE)

bullets(s5, 0.4, 3.6, 5.9, 1.5, [
    ('Demographics', 16, True, C_DARK_BLUE, False),
    ('Gender, Senior Citizen, Partner, Dependents', 13, False, C_GRAY, True),
    ('Service Subscriptions', 16, True, C_DARK_BLUE, False),
    ('Phone, Internet, MultipleLines, Add-ons (x6)', 13, False, C_GRAY, True),
])
bullets(s5, 6.6, 3.6, 6.2, 1.5, [
    ('Contract & Billing', 16, True, C_DARK_BLUE, False),
    ('Contract Type, Paperless Billing, Payment Method', 13, False, C_GRAY, True),
    ('Charges', 16, True, C_DARK_BLUE, False),
    ('MonthlyCharges, TotalCharges, Tenure (months)', 13, False, C_GRAY, True),
])

txb(s5, 0.4, 5.25, 12.5, 0.38, 'Stratified 3-Way Split',
    size=15, bold=True, color=C_DARK_BLUE)
add_table(s5, 0.4, 5.7, 7.0,
    [['Split',      'Proportion', 'Size',  'Purpose'],
     ['Train',      '70%',        '4,930', 'Model training + CV'],
     ['Validation', '15%',        '1,057', 'Calibration + threshold tuning'],
     ['Test',       '15%',        '1,057', 'Final held-out evaluation']],
    [2.0, 1.5, 1.5, 2.5], font_size=12)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 6 — ML PIPELINE ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════════════
s6 = new_slide('ML Pipeline Architecture')

stages = [
    ('Raw CSV',           C_MID_BLUE),
    ('Data Validation',   C_MID_BLUE),
    ('Feature Eng.',      C_MID_BLUE),
    ('3-Way Split',       C_MID_BLUE),
    ('Baseline Ladder',   C_MID_BLUE),
    ('Optuna Tuning',     C_ORANGE),
    ('Ensembling',        C_ORANGE),
    ('Calibration',       C_PURPLE),
    ('Threshold Select',  C_PURPLE),
    ('Evaluation',        C_GREEN),
    ('SHAP Explain.',     C_GREEN),
    ('Deploy Chatbot',    C_GREEN),
]

row1 = stages[:6]
row2 = list(reversed(stages[6:]))   # Deploy, SHAP, Eval, Threshold, Cal, Ensembling

bw, bh  = 1.88, 0.8
gap     = 0.095
row1_y  = 2.05
row2_y  = 3.25

# ── Row 1 (left to right) ────────────────────────────────────────────────────
for i, (label, color) in enumerate(row1):
    lft = 0.35 + i * (bw + gap)
    rect = s6.shapes.add_shape(9, Inches(lft), Inches(row1_y), Inches(bw), Inches(bh))
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.color.rgb = C_WHITE
    tb_r = s6.shapes.add_textbox(Inches(lft), Inches(row1_y), Inches(bw), Inches(bh))
    tf   = tb_r.text_frame
    tf.word_wrap = True
    p    = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r    = p.add_run()
    r.text           = label
    r.font.size      = Pt(11)
    r.font.bold      = True
    r.font.color.rgb = C_WHITE
    # Right-pointing arrow between blocks (except after last block)
    if i < len(row1) - 1:
        arr = s6.shapes.add_textbox(
            Inches(lft + bw), Inches(row1_y + 0.2), Inches(gap), Inches(0.4))
        arr.text_frame.paragraphs[0].add_run().text = '>'
        arr.text_frame.paragraphs[0].runs[0].font.color.rgb = C_WHITE
        arr.text_frame.paragraphs[0].runs[0].font.size      = Pt(10)
        arr.text_frame.paragraphs[0].runs[0].font.bold      = True

# ── Turn arrow (right side, between row1 and row2) ───────────────────────────
# Last block of row1 ends at: 0.35 + 5*(bw+gap) + bw = 0.35 + 9.875 + 1.88 = 12.105
# Place turn arrow to the right of that block, centred vertically between the rows
last_right = 0.35 + 5 * (bw + gap) + bw          # = 12.105"
turn_x     = last_right + 0.08                     # small gap after last block

turn = s6.shapes.add_textbox(
    Inches(turn_x), Inches(row1_y + 0.15), Inches(0.5), Inches(row2_y - row1_y + bh * 0.7))
p_t  = turn.text_frame.paragraphs[0]
r_t  = p_t.add_run()
r_t.text           = 'v'
r_t.font.size      = Pt(18)
r_t.font.bold      = True
r_t.font.color.rgb = C_ORANGE

# ── Row 2 (flow: right to left, so Ensembling is rightmost) ──────────────────
for i, (label, color) in enumerate(row2):
    lft = 0.35 + i * (bw + gap)
    rect = s6.shapes.add_shape(9, Inches(lft), Inches(row2_y), Inches(bw), Inches(bh))
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.color.rgb = C_WHITE
    tb_r = s6.shapes.add_textbox(Inches(lft), Inches(row2_y), Inches(bw), Inches(bh))
    tf   = tb_r.text_frame
    tf.word_wrap = True
    p    = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r    = p.add_run()
    r.text           = label
    r.font.size      = Pt(11)
    r.font.bold      = True
    r.font.color.rgb = C_WHITE
    # Left-pointing arrow between blocks in row2 (flow is right to left)
    if i < len(row2) - 1:
        arr = s6.shapes.add_textbox(
            Inches(lft + bw), Inches(row2_y + 0.2), Inches(gap), Inches(0.4))
        arr.text_frame.paragraphs[0].add_run().text = '<'
        arr.text_frame.paragraphs[0].runs[0].font.color.rgb = C_WHITE
        arr.text_frame.paragraphs[0].runs[0].font.size      = Pt(10)
        arr.text_frame.paragraphs[0].runs[0].font.bold      = True

# ── Legend ────────────────────────────────────────────────────────────────────
txb(s6, 0.4, 4.25, 1.2, 0.35, 'Legend:', size=12, bold=True, color=C_GRAY)
for i, (lbl, clr) in enumerate([
        ('Data Prep', C_MID_BLUE),
        ('Modeling', C_ORANGE),
        ('Eval / Deploy', C_GREEN),
        ('Calibration', C_PURPLE)]):
    rx = 1.2 + i * 2.9
    rect = s6.shapes.add_shape(1, Inches(rx), Inches(4.62), Inches(0.22), Inches(0.22))
    rect.fill.solid()
    rect.fill.fore_color.rgb = clr
    rect.line.fill.background()
    txb(s6, rx + 0.28, 4.57, 2.4, 0.35, lbl, size=12, color=C_GRAY)

# ── Key notes ─────────────────────────────────────────────────────────────────
bullets(s6, 0.4, 5.1, 12.2, 1.6, [
    ('Zero data leakage: SMOTE applied inside CV folds only - never on held-out data', 13, False, C_GRAY, True),
    ('Optuna 200 trials, TPE Sampler + Median Pruner: Bayesian hyperparameter optimization', 13, False, C_GRAY, True),
    ('Threshold 0.22 selected via PR-curve max-F1 with business Expected Value verification', 13, False, C_GRAY, True),
])


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 7 — FEATURE ENGINEERING
# ══════════════════════════════════════════════════════════════════════════════
s7 = new_slide('Feature Engineering  -  12 Domain-Driven Features')

txb(s7, 0.4, 1.9, 12.5, 0.35,
    'Engineered from raw data to capture behavioral and risk signals:',
    size=13, italic=True, color=C_GRAY)

add_table(s7, 0.4, 2.3, 12.53,
    [['Feature',               'Description',                                        'Type'],
     ['TenureGroup',           'Bucketed tenure: 0-6m, 6-12m, 12-24m, 24-48m, 48m+','Categorical'],
     ['IsNewCustomer',         'Binary flag: tenure <= 6 months = 1',                'Binary'],
     ['AvgMonthlySpend',       'TotalCharges / (tenure + 1)',                         'Numeric'],
     ['SpendRatio',            'MonthlyCharges / (AvgMonthlySpend + e)',              'Numeric'],
     ['SpendAcceleration',     'MonthlyCharges - AvgMonthlySpend',                   'Numeric'],
     ['NumServices',           'Count of active add-on services (max 6)',             'Numeric'],
     ['HasProtectionBundle',   'OnlineSecurity AND DeviceProtection = 1',             'Binary'],
     ['HasStreamingBundle',    'StreamingTV AND StreamingMovies = 1',                 'Binary'],
     ['ContractRiskScore',     'Month-to-month=3, One year=1, Two year=0',            'Ordinal'],
     ['HighRiskPayment',       'Electronic check payment = 1',                        'Binary'],
     ['FamilyStability',       'Partner OR Dependents = 1',                           'Binary'],
     ['IsFiberCustomer',       'Fiber optic internet = 1',                            'Binary']],
    [3.8, 6.5, 2.0], font_size=11)

txb(s7, 0.4, 6.9, 12.5, 0.35,
    'Note: engineer_features() is identically replicated in both training pipeline and'
    ' chatbot to ensure train-serve feature parity.',
    size=11, italic=True, color=C_ORANGE)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 8 — HANDLING CLASS IMBALANCE
# ══════════════════════════════════════════════════════════════════════════════
s8 = new_slide('Handling Class Imbalance  (73:27 Ratio)')

strategies = [
    ('1.  SMOTE', C_MID_BLUE,
     '- Synthetic Minority Oversampling\n- Applied INSIDE ImbPipeline\n'
     '- Only within CV folds\n- Zero leakage to validation data'),
    ('2.  Class Weights', C_ORANGE,
     '- balanced class_weight on all estimators\n'
     '- Penalizes minority misclassification\n- Works alongside SMOTE'),
    ('3.  Threshold Optimization', C_GREEN,
     '- Default threshold = 0.50\n- Production threshold = 0.22\n'
     '- Selected via PR-curve max-F1\n- Reflects 10:1 cost asymmetry'),
]
bw_s, bh_s = 3.8, 2.8
for i, (title, color, body) in enumerate(strategies):
    lft = 0.4 + i * (bw_s + 0.38)
    rect = s8.shapes.add_shape(9, Inches(lft), Inches(2.0), Inches(bw_s), Inches(bh_s))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0xF0, 0xF7, 0xFF)
    rect.line.color.rgb = color
    hdr = s8.shapes.add_shape(1, Inches(lft), Inches(2.0), Inches(bw_s), Inches(0.5))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = color
    hdr.line.fill.background()
    tb_h = s8.shapes.add_textbox(Inches(lft + 0.05), Inches(2.0), Inches(bw_s - 0.1), Inches(0.5))
    p    = tb_h.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r    = p.add_run()
    r.text           = title
    r.font.size      = Pt(14)
    r.font.bold      = True
    r.font.color.rgb = C_WHITE
    tb_b = s8.shapes.add_textbox(Inches(lft + 0.1), Inches(2.6), Inches(bw_s - 0.2), Inches(bh_s - 0.75))
    tb_b.text_frame.word_wrap = True
    p2   = tb_b.text_frame.paragraphs[0]
    r2   = p2.add_run()
    r2.text           = body
    r2.font.size      = Pt(12)
    r2.font.color.rgb = C_GRAY

txb(s8, 0.4, 5.0, 12.5, 0.35, 'Combined Effect:', size=14, bold=True, color=C_DARK_BLUE)
bullets(s8, 0.4, 5.4, 12.2, 1.5, [
    ('All three strategies together provide complementary imbalance correction', 13, False, C_GRAY, True),
    ('SMOTE: synthetic sampling | Class weights: loss function | Threshold: decision boundary', 13, False, C_GRAY, True),
    ('Final threshold 0.22 maximizes Expected Value under 10:1 cost asymmetry', 13, False, C_GRAY, True),
])


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 9 — MODEL DEVELOPMENT
# ══════════════════════════════════════════════════════════════════════════════
s9 = new_slide('Model Development  -  Optuna Hyperparameter Optimization')

txb(s9, 0.4, 1.9, 5.9, 0.38, 'Baseline Ladder', size=15, bold=True, color=C_DARK_BLUE)
bullets(s9, 0.4, 2.3, 5.7, 1.5, [
    ('Logistic Regression (LR)  - linear baseline',       13, False, C_GRAY, True),
    ('Random Forest (RF)  - ensemble baseline',            13, False, C_GRAY, True),
    ('XGBoost (XGB)  - gradient boosting',                 13, False, C_GRAY, True),
    ('LightGBM (LGBM)  - histogram-based gradient boost',  13, False, C_GRAY, True),
    ('CatBoost  - categorical-native boosting',            13, False, C_GRAY, True),
])

txb(s9, 0.4, 3.9, 5.9, 0.38, 'Cross-Validation Strategy', size=15, bold=True, color=C_DARK_BLUE)
bullets(s9, 0.4, 4.32, 5.7, 1.5, [
    ('RepeatedStratifiedKFold: 5 folds x 2 repeats = 10 evaluations', 13, False, C_GRAY, True),
    ('Lower variance than standard K-Fold',                13, False, C_GRAY, True),
    ('Stratification preserves 73:27 class ratio in every fold', 13, False, C_GRAY, True),
])

add_table(s9, 6.6, 1.9, 6.0,
    [['Optuna Setting',  'Value'],
     ['Sampler',         'TPE (Bayesian search)'],
     ['Pruner',          'Median Pruner'],
     ['Trials',          '200 per model'],
     ['Objective',       'Average Precision (AP)'],
     ['CV',              '5 folds x 2 repeats'],
     ['SMOTE',           'Inside CV folds only'],
     ['Importance',      'fANOVA ranking']],
    [3.0, 3.0], font_size=12)

txb(s9, 6.6, 5.1, 6.0, 0.35, 'Why Average Precision?', size=13, bold=True, color=C_DARK_BLUE)
txb(s9, 6.6, 5.5, 6.0, 1.0,
    'AP is more sensitive than ROC-AUC under class imbalance. It focuses on '
    'the minority (churn) class across all thresholds.',
    size=12, color=C_GRAY, wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 10 — ENSEMBLE & CALIBRATION
# ══════════════════════════════════════════════════════════════════════════════
s10 = new_slide('Ensemble Methods & Probability Calibration')

txb(s10, 0.4, 1.9, 12.5, 0.38, 'Three Ensemble Strategies Evaluated',
    size=15, bold=True, color=C_DARK_BLUE)

ens_data = [
    ('Soft Voting',  C_MID_BLUE, 'Probability averaging across all\nbase models. Simple, robust,\nlow variance.'),
    ('Single Best',  C_ORANGE,   'Top individual model from Optuna.\nLightest inference.\nHighest Recall.'),
    ('Stacking',     C_PURPLE,   'Meta-learner (Logistic Regression)\ntrained on base-model\nout-of-fold predictions.'),
]
for i, (name, color, desc) in enumerate(ens_data):
    lft = 0.4 + i * 4.1
    rect = s10.shapes.add_shape(9, Inches(lft), Inches(2.35), Inches(3.9), Inches(1.6))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0xF0, 0xF7, 0xFF)
    rect.line.color.rgb = color
    hdr = s10.shapes.add_shape(1, Inches(lft), Inches(2.35), Inches(3.9), Inches(0.42))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = color
    hdr.line.fill.background()
    tb_h = s10.shapes.add_textbox(Inches(lft + 0.05), Inches(2.35), Inches(3.8), Inches(0.42))
    p    = tb_h.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r    = p.add_run()
    r.text           = name
    r.font.size      = Pt(14)
    r.font.bold      = True
    r.font.color.rgb = C_WHITE
    tb_b = s10.shapes.add_textbox(Inches(lft + 0.1), Inches(2.82), Inches(3.7), Inches(1.1))
    tb_b.text_frame.word_wrap = True
    p2   = tb_b.text_frame.paragraphs[0]
    r2   = p2.add_run()
    r2.text           = desc
    r2.font.size      = Pt(12)
    r2.font.color.rgb = C_GRAY

txb(s10, 0.4, 4.1, 12.5, 0.38, 'Probability Calibration - Why It Matters',
    size=15, bold=True, color=C_DARK_BLUE)
bullets(s10, 0.4, 4.52, 12.2, 1.6, [
    ('Method: Isotonic Regression on held-out validation set', 13, False, C_GRAY, True),
    ('Ensures predicted 70% probability = 70% empirical churn rate (reliable for business)', 13, False, C_GRAY, True),
    ('Brier Score improvement tracked; reliability diagrams generated pre/post calibration', 13, False, C_GRAY, True),
    ('Platt Scaling tested but Isotonic performs better for n > 1,000 samples', 13, False, C_GRAY, True),
])

txb(s10, 0.4, 6.3, 12.5, 0.5,
    'WINNER: Single Best (Calibrated) - highest Recall (0.825) + F1 (0.636) + EV per user ($0.155)',
    size=14, bold=True, color=C_GREEN)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 11 — RESULTS & MODEL COMPARISON
# ══════════════════════════════════════════════════════════════════════════════
s11 = new_slide('Results & Model Comparison')

txb(s11, 0.4, 1.9, 12.5, 0.35,
    'Calibrated Model Performance on Test Set  (Threshold = 0.22)',
    size=14, bold=True, color=C_DARK_BLUE)

add_table(s11, 0.4, 2.28, 12.53,
    [['Model',                 'ROC-AUC', 'Avg Precision', 'Recall',   'F1',     'MCC',  'EV/User'],
     ['* Single Best (Cal.)',  '0.845',   '0.642',         '0.825 *',  '0.636 *','0.490','$0.155 *'],
     ['Soft Voting (Cal.)',    '0.847',   '-',             '0.764',    '0.636',  '-',    '$0.137'],
     ['Stacking (Cal.)',       '0.850 *', '-',             '0.779',    '0.636',  '-',    '$0.141']],
    [3.5, 1.5, 2.0, 1.4, 1.2, 1.2, 1.5], font_size=12)

txb(s11, 0.4, 3.9, 12.5, 0.38, 'Final Production Model Scorecard',
    size=14, bold=True, color=C_DARK_BLUE)

scorecard = [
    ('ROC-AUC',      '0.8503'),
    ('Avg Precision','0.6419'),
    ('Recall',       '0.825'),
    ('F1 Score',     '0.636'),
    ('MCC',          '0.490'),
    ('Brier Score',  '0.134'),
    ('EV / User',    '$0.15'),
    ('Total EV',     '$164'),
]
sw = 2.88
for i, (lbl, val) in enumerate(scorecard):
    row = i // 4
    col = i % 4
    info_box(s11, 0.4 + col * (sw + 0.24), 4.38 + row * 1.22, sw, 1.05,
             lbl, val, bg=C_LIGHT_BLUE)

txb(s11, 0.4, 6.88, 12.5, 0.38,
    'Production Threshold = 0.22  |  Test Set = 1,057 customers',
    size=12, italic=True, color=C_MID_BLUE)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 12 — SHAP EXPLAINABILITY
# ══════════════════════════════════════════════════════════════════════════════
s12 = new_slide('SHAP Explainability')

txb(s12, 0.4, 1.9, 5.9, 0.38, 'Top 5 Churn Drivers (SHAP)',
    size=15, bold=True, color=C_DARK_BLUE)

drivers = [
    ('ContractRiskScore',  'Month-to-month churns 3x vs 2-year contracts', C_RED),
    ('IsNewCustomer',      'Tenure <= 6 months = highest individual risk',  C_ORANGE),
    ('HighRiskPayment',    'Electronic check strongly correlates w/ churn', C_ORANGE),
    ('IsFiberCustomer',    'Fiber customers churn despite premium pricing',  C_MID_BLUE),
    ('NumServices',        'Low add-on engagement = higher churn risk',     C_MID_BLUE),
]
for i, (feat, desc, color) in enumerate(drivers):
    rect = s12.shapes.add_shape(1, Inches(0.4), Inches(2.35 + i * 0.82), Inches(5.8), Inches(0.72))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0xF2, 0xF2, 0xF2)
    rect.line.color.rgb = color
    bar  = s12.shapes.add_shape(1, Inches(0.4), Inches(2.35 + i * 0.82), Inches(0.15), Inches(0.72))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    txb(s12, 0.65, 2.37 + i * 0.82, 5.45, 0.35,
        f'#{i+1}  {feat}', size=12, bold=True, color=C_DARK_GRAY)
    txb(s12, 0.65, 2.72 + i * 0.82, 5.45, 0.32,
        desc, size=11, color=C_GRAY)

txb(s12, 6.6, 1.9, 6.0, 0.38, 'Analysis Types',
    size=15, bold=True, color=C_DARK_BLUE)

shap_types = [
    ('Global Explainability',    'Feature importance bar chart\nBeeswarm plot across all predictions'),
    ('Local Explainability',     'Waterfall plot per customer\nForce plot visualization'),
    ('Dependence Plots',         'Non-linear feature effects\nFeature interaction analysis'),
    ('Importance Triangulation', 'Built-in (gain) vs Permutation\nvs SHAP - cross-validated'),
]
for i, (title, desc) in enumerate(shap_types):
    txb(s12, 6.6, 2.35 + i * 1.12, 6.0, 0.38, title, size=13, bold=True, color=C_MID_BLUE)
    txb(s12, 6.6, 2.72 + i * 1.12, 6.0, 0.65, desc,  size=12, color=C_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 13 — DEPLOYMENT / CHATBOT
# ══════════════════════════════════════════════════════════════════════════════
s13 = new_slide('Deployment  -  AI-Powered Streamlit Chatbot')

txb(s13, 0.4, 1.9, 5.9, 0.38, 'Chatbot Features',
    size=15, bold=True, color=C_DARK_BLUE)
bullets(s13, 0.4, 2.3, 5.8, 3.2, [
    ('19 guided questions - one at a time', 13, False, C_GRAY, True),
    ('Clickable option buttons + numeric inputs', 13, False, C_GRAY, True),
    ('Real-time progress bar', 13, False, C_GRAY, True),
    ('Color-coded result card (Churn / No-Churn)', 13, False, C_GRAY, True),
    ('AI explanation via Claude API (4-5 sentences)', 13, False, C_GRAY, True),
    ('Customer summary table', 13, False, C_GRAY, True),
    ('Multiple predictions per session', 13, False, C_GRAY, True),
    ('Rule-based fallback if API unavailable', 13, False, C_GRAY, True),
])

txb(s13, 6.6, 1.9, 6.0, 0.38, 'Inference Pipeline',
    size=15, bold=True, color=C_DARK_BLUE)

pipeline_steps = [
    ('User Input',          '19 guided questions',       C_MID_BLUE),
    ('Feature Engineering', 'engineer_features() x 12',  C_MID_BLUE),
    ('ML Model',            'predict_proba() @ 0.22',    C_ORANGE),
    ('Claude API',          'Explanation generation',    C_PURPLE),
    ('Output',              'Probability + Explanation', C_GREEN),
]
for i, (step, detail, color) in enumerate(pipeline_steps):
    rect = s13.shapes.add_shape(9, Inches(6.6), Inches(2.35 + i * 0.92), Inches(5.8), Inches(0.75))
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.color.rgb = C_WHITE
    txb(s13, 6.65, 2.37 + i * 0.92, 2.8,  0.38, step,   size=12, bold=True,  color=C_WHITE)
    txb(s13, 9.5,  2.37 + i * 0.92, 2.85, 0.38, detail, size=11, bold=False, color=C_WHITE)
    if i < len(pipeline_steps) - 1:
        arr = s13.shapes.add_textbox(
            Inches(9.35), Inches(3.06 + i * 0.92), Inches(0.5), Inches(0.35))
        arr.text_frame.paragraphs[0].add_run().text = 'v'
        arr.text_frame.paragraphs[0].runs[0].font.color.rgb = color
        arr.text_frame.paragraphs[0].runs[0].font.size      = Pt(11)
        arr.text_frame.paragraphs[0].runs[0].font.bold      = True

txb(s13, 0.4, 5.65, 12.5, 0.38, 'Deployment Artifacts',
    size=14, bold=True, color=C_DARK_BLUE)
add_table(s13, 0.4, 6.08, 12.53,
    [['Artifact',          'Description'],
     ['churn_model.pkl',   'Serialized pipeline via joblib'],
     ['model_card.json',   'Metrics, threshold, feature list, training date'],
     ['Drift Monitor',     'Feature distribution tracking for production'],
     ['chatbot.py',        'Streamlit app - run: streamlit run chatbot.py']],
    [4.0, 8.5], font_size=11)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 14 — CONCLUSION & FUTURE WORK
# ══════════════════════════════════════════════════════════════════════════════
s14 = new_slide('Conclusion & Future Work')

txb(s14, 0.4, 1.9, 12.5, 0.38, 'Key Contributions',
    size=15, bold=True, color=C_DARK_BLUE)
bullets(s14, 0.4, 2.3, 12.2, 2.5, [
    ('Production-grade pipeline: zero data leakage, SMOTE inside CV folds only', 14, False, C_GRAY, True),
    ('Optuna Bayesian optimization: 200 trials, Avg Precision objective, TPE + MedianPruner', 14, False, C_GRAY, True),
    ('Business-cost-aware design: 10:1 cost ratio -> threshold 0.22 -> Recall = 0.825', 14, False, C_GRAY, True),
    ('Isotonic calibration: statistically reliable probability estimates', 14, False, C_GRAY, True),
    ('SHAP triangulated explainability: global + local + dependence + importance triangulation', 14, False, C_GRAY, True),
    ('AI-powered chatbot: conversational interface with Claude API explanations', 14, False, C_GRAY, True),
])

txb(s14, 0.4, 4.95, 12.5, 0.38, 'Future Work',
    size=15, bold=True, color=C_DARK_BLUE)
bullets(s14, 0.4, 5.38, 12.2, 1.3, [
    ('Online learning for concept drift adaptation in production', 13, False, C_GRAY, True),
    ('Richer behavioral features: call logs, support tickets, app usage', 13, False, C_GRAY, True),
    ('Multi-objective optimization incorporating fairness constraints', 13, False, C_GRAY, True),
])

for i, (lbl, val) in enumerate([('ROC-AUC', '0.85'), ('Recall', '0.825'),
                                  ('F1', '0.636'), ('EV Total', '$164')]):
    info_box(s14, 0.4 + i * 3.1, 6.7, 2.85, 0.75, lbl, val, bg=C_LIGHT_BLUE)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 15 — THANK YOU
#  Uses the existing blank slide (currently at index 2).
#  After adding s3-s14, this slide sits at index 2 — we move it to the END.
# ══════════════════════════════════════════════════════════════════════════════
s15 = prs.slides[2]   # existing blank slide

rect_bg = s15.shapes.add_shape(1, Inches(2.0), Inches(2.3), Inches(9.33), Inches(2.2))
rect_bg.fill.solid()
rect_bg.fill.fore_color.rgb = C_DARK_BLUE
rect_bg.line.fill.background()

txb(s15, 2.0, 2.35, 9.33, 1.0, 'Thank You!',
    size=40, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
txb(s15, 2.0, 3.35, 9.33, 0.6,
    'Questions & Discussion',
    size=20, color=C_LIGHT_BLUE, align=PP_ALIGN.CENTER)

txb(s15, 0.4, 5.0, 12.5, 0.38, 'References',
    size=15, bold=True, color=C_DARK_BLUE)
txb(s15, 0.4, 5.45, 12.5, 1.8,
    '[1] Akiba et al., "Optuna," KDD 2019.   '
    '[2] Lundberg & Lee, "SHAP," NeurIPS 2017.   '
    '[3] Chawla et al., "SMOTE," JAIR 2002.   '
    '[4] Chen & Guestrin, "XGBoost," KDD 2016.   '
    '[5] Ke et al., "LightGBM," NeurIPS 2017.   '
    '[6] Reichheld & Schefter, "E-Loyalty," HBR 2000.',
    size=11, color=C_GRAY, wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
#  FIX SLIDE ORDER: move Thank You from index 2 to the last position
#  Without this fix, the deck order would be:
#    [Title(0), Outline(1), ThankYou(2), Intro(3), ..., Conclusion(14)]
#  After the move:
#    [Title(0), Outline(1), Intro(2), ..., Conclusion(13), ThankYou(14)]
# ══════════════════════════════════════════════════════════════════════════════
move_slide(prs, old_idx=2, new_idx=len(prs.slides) - 1)


# ── Save ──────────────────────────────────────────────────────────────────────
prs.save(OUTPUT)
print('[OK] Saved: ' + OUTPUT)
print('     Total slides: ' + str(len(prs.slides)))
print('     Slide order:')
for i, sl in enumerate(prs.slides):
    texts = [sh.text_frame.text[:40] for sh in sl.shapes if sh.has_text_frame]
    print(f'     {i+1:2d}. {texts[0] if texts else "(no text)"}')
